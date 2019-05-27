# -*- coding:UTF-8 -*-
from pptx import Presentation
import trans

def trans_ppt(source_ppt,out_ppt):
    prs = Presentation(source_ppt)
    for slide in prs.slides:
        for index,shape in enumerate(slide.shapes):
            # print("index : " + str(index))
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    print(run.text)
                    translation = trans.baidu.tochinese(run.text)
                    paragraph.text = translation
                    print(translation)
    prs.save(out_ppt)


trans_ppt("resource/2019 TM kick off deck_V3.0.pptx","aaa.pptx")